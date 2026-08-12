"""
PowerPoint tool — create and edit PPTX files.
Uses python-pptx.

Usage:
  python scripts/pptx_tool.py create output.pptx
  python scripts/pptx_tool.py edit input.pptx output.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def create(path: str, title: str = "Название"):
    """Create a new PPTX with a title slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    prs.save(path)
    print(f"Created: {path}")


def edit(input_path: str, output_path: str, commands: list[dict] | None = None):
    """Edit an existing PPTX."""
    prs = Presentation(input_path)

    if commands:
        for cmd in commands:
            _run_command(prs, cmd)

    prs.save(output_path)
    print(f"Saved: {output_path}")


def _run_command(prs: Presentation, cmd: dict):
    kind = cmd.get("type")

    if kind == "add_slide":
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if "text" in cmd:
            txBox = slide.shapes.add_textbox(
                Inches(cmd.get("left", 1)),
                Inches(cmd.get("top", 0.5)),
                Inches(cmd.get("width", 11)),
                Inches(cmd.get("height", 1)),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cmd["text"]
            p.font.size = Pt(cmd.get("font_size", 24))

    elif kind == "add_table":
        rows, cols = cmd["rows"], cmd["cols"]
        left = Inches(cmd.get("left", 1))
        top = Inches(cmd.get("top", 1))
        width = Inches(cmd.get("width", 11))
        height = Inches(cmd.get("height", 0.5 * rows))
        table_shape = prs.slides[-1].shapes.add_table(rows, cols, left, top, width, height)
        for ri, row_data in enumerate(cmd.get("data", [])):
            for ci, cell_text in enumerate(row_data):
                cell = table_shape.table.cell(ri, ci)
                cell.text = str(cell_text)

    elif kind == "set_title":
        slide = prs.slides[-1]
        txBox = slide.shapes.add_textbox(
            Inches(cmd.get("left", 1)),
            Inches(cmd.get("top", 0.3)),
            Inches(cmd.get("width", 11)),
            Inches(cmd.get("height", 1)),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = cmd["text"]
        p.font.size = Pt(cmd.get("font_size", 36))
        p.font.bold = True


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    command = sys.argv[1]

    if command == "create":
        out = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
        title = sys.argv[3] if len(sys.argv) > 3 else "Новый документ"
        create(out, title)

    elif command == "edit":
        if len(sys.argv) < 3:
            print("Need input and output paths")
            sys.exit(1)
        inp = sys.argv[2]
        out = sys.argv[3] if len(sys.argv) > 3 else inp
        edit(inp, out)

    else:
        print(f"Unknown command: {command}")
        sys.exit(1)


if __name__ == "__main__":
    main()
