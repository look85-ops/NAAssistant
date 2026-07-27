#!/usr/bin/env python3
"""Утилита для чтения/создания/правки файлов планирования.
Поддерживает: xlsx, docx, pptx, pdf."""

import argparse
import sys
import io
from pathlib import Path
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font, Alignment, PatternFill

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")


HEADER_FILL = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
HEADER_FONT = Font(color="FFFFFF", bold=True, size=11)


# ── Определитель типа файла ──────────────────────────────────────

def _ext(file):
    return Path(file).suffix.lower()


# ── READ ──────────────────────────────────────────────────────────

def _read_xlsx(file):
    wb = load_workbook(file, data_only=True)
    ws = wb.active
    print(f"Лист: {ws.title}  (строк: {ws.max_row}, колонок: {ws.max_column})\n")
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=True):
        vals = [str(v) if v is not None else "" for v in row]
        print("\t".join(vals))
    wb.close()


def _read_docx(file):
    from docx import Document
    doc = Document(file)
    print(f"Абзацев: {len(doc.paragraphs)}, таблиц: {len(doc.tables)}\n")
    for p in doc.paragraphs:
        if p.text.strip():
            print(p.text)
    for ti, t in enumerate(doc.tables, 1):
        print(f"\n--- Таблица {ti} ---")
        for row in t.rows:
            print("\t".join(c.text for c in row.cells))


def _read_pptx(file):
    from pptx import Presentation
    prs = Presentation(file)
    print(f"Слайдов: {len(prs.slides)}\n")
    for si, slide in enumerate(prs.slides, 1):
        print(f"--- Слайд {si} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        print(p.text)
            if shape.has_table:
                for row in shape.table.rows:
                    print("\t".join(c.text for c in row.cells))


def _read_pdf(file):
    import fitz
    doc = fitz.open(file)
    print(f"Страниц: {doc.page_count}\n")
    for page in doc:
        print(page.get_text())
        print("---")
    doc.close()


def cmd_read(args):
    fpath = args.file
    print(f"Файл: {fpath}")
    ext = _ext(fpath)
    readers = {
        ".xlsx": _read_xlsx,
        ".xls": _read_xlsx,
        ".docx": _read_docx,
        ".pptx": _read_pptx,
        ".pdf": _read_pdf,
    }
    handler = readers.get(ext)
    if not handler:
        print(f"Формат {ext} не поддерживается. Поддерживаются: xlsx, docx, pptx, pdf")
        sys.exit(1)
    handler(fpath)


# ── CREATE (только xlsx) ─────────────────────────────────────────

def _write_sheet(ws, headers, rows):
    for i, h in enumerate(headers, 1):
        c = ws.cell(row=1, column=i, value=h)
        c.font = HEADER_FONT
        c.fill = HEADER_FILL
        c.alignment = Alignment(horizontal="center")
    for r_idx, row in enumerate(rows, 2):
        for c_idx, val in enumerate(row, 1):
            ws.cell(row=r_idx, column=c_idx, value=val)
    ws.column_dimensions["A"].width = 24
    for col in range(2, len(headers) + 1):
        ws.column_dimensions[chr(64 + col) if col <= 26 else "Z"].width = 18


TEMPLATES = {
    "budget": {
        "headers": ["Категория", "План (BYN)", "Факт (BYN)", "Разница", "Примечание"],
        "rows": [
            ["Аренда", "", "", "", ""],
            ["Коммунальные", "", "", "", ""],
            ["Продукты", "", "", "", ""],
            ["Транспорт", "", "", "", ""],
            ["Связь/Интернет", "", "", "", ""],
            ["Здоровье/Страховка", "", "", "", ""],
            ["Кафе/Досуг", "", "", "", ""],
            ["Одежда/Быт", "", "", "", ""],
            ["Непредвиденное (10%)", "", "", "", ""],
            ["ИТОГО", "", "", "", ""],
        ],
    },
    "timeline": {
        "headers": ["Месяц", "Событие", "Дедлайн", "Статус", "Бюджет (BYN)"],
        "rows": [
            ["Июль 2026", "", "", "", ""],
            ["Август 2026", "", "", "", ""],
            ["Сентябрь 2026", "", "", "", ""],
            ["Октябрь 2026", "", "", "", ""],
            ["Ноябрь 2026", "", "", "", ""],
            ["Декабрь 2026", "", "", "", ""],
        ],
    },
    "task": {
        "headers": ["ID", "Задача", "Категория", "Приоритет", "Статус", "Дедлайн"],
        "rows": [],
    },
}


def cmd_create(args):
    if Path(args.file).exists():
        print(f"Ошибка: {args.file} уже существует.")
        sys.exit(1)
    tpl = TEMPLATES.get(args.type, TEMPLATES["budget"])
    wb = Workbook()
    ws = wb.active
    ws.title = args.type.capitalize()
    _write_sheet(ws, tpl["headers"], tpl["rows"])
    wb.save(args.file)
    wb.close()
    print(f"Создан: {args.file}  (тип: {args.type})")


# ── UPDATE (только xlsx) ─────────────────────────────────────────

def cmd_update(args):
    wb = load_workbook(args.file)
    ws = wb.active
    ws.cell(row=args.row, column=args.col, value=args.value)
    wb.save(args.file)
    wb.close()
    print(f"Обновлено [{args.row},{args.col}] = {args.value}")


# ── LIST ──────────────────────────────────────────────────────────

def cmd_list(args):
    exts = [".xlsx", ".docx", ".pptx", ".pdf"]
    pattern = args.pattern or "*"
    for f in sorted(Path(args.dir).glob(pattern)):
        if f.suffix.lower() in exts:
            print(f"{f.name}  ({_fmt_size(f.stat().st_size)})")


def _fmt_size(b):
    for u in ("B", "KB", "MB"):
        if b < 1024:
            return f"{b:.0f} {u}"
        b /= 1024
    return f"{b:.1f} GB"


# ── CLI ───────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Планировщик: читает xlsx/docx/pptx/pdf, создаёт/правит xlsx")
    sub = parser.add_subparsers(dest="command")

    p_read = sub.add_parser("read", help="Прочитать файл (xlsx/docx/pptx/pdf)")
    p_read.add_argument("file")

    p_create = sub.add_parser("create", help="Создать xlsx-шаблон")
    p_create.add_argument("file")
    p_create.add_argument("--type", choices=list(TEMPLATES), default="budget")

    p_update = sub.add_parser("update", help="Обновить ячейку в xlsx")
    p_update.add_argument("file")
    p_update.add_argument("--row", type=int, required=True)
    p_update.add_argument("--col", type=int, required=True)
    p_update.add_argument("--value", required=True)

    p_list = sub.add_parser("list", help="Список файлов в директории")
    p_list.add_argument("--dir", default=".")
    p_list.add_argument("--pattern", default="*")

    args = parser.parse_args()
    if not args.command:
        parser.print_help()
        sys.exit(1)

    cmds = {
        "read": cmd_read,
        "create": cmd_create,
        "update": cmd_update,
        "list": cmd_list,
    }
    cmds[args.command](args)


if __name__ == "__main__":
    main()
